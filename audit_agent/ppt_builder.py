"""
Exportación del informe a PowerPoint sobre la plantilla corporativa
(`config/plantilla_informe.pptx`, saneada con `scripts/sanear_plantilla.py`).

No se dibuja nada desde cero: cada diapositiva del informe es una diapositiva
de la plantilla (o un duplicado de ella) en la que se sustituye el texto
conservando fuentes, tamaños, colores, tablas e imágenes. Estructura:

  1. Portada            (plantilla d1: nombre, lista de distribución, fecha y ref.)
  2. Índice             (d2, tal cual)
  3. Introducción       (d3: frase fija del plan, bloques Contexto / Objetivo /
                         Riesgos a cubrir / Alcance / … y frase de normas al pie)
  4. Portadilla         (d4) + Resumen ejecutivo (d5: texto, punto de la escala de
                         Evaluación Global y Próximos pasos)
  5. Portadilla         (d6) + una diapositiva de tabla por recomendación (duplicado
                         de d7: banda RIESGO, «NN Título» + prosa, caja gris de
                         detalles, consecuencias, Recomendación N.k, Área /
                         Responsable / Plazo); si no cabe, «(continuación)»
  6. Portadilla         (d9) + una por sugerencia (duplicado de d10)
  7. Portadilla         (d11) + Anexo de planes de acción

Reglas python-pptx sobre plantillas: nunca asignar `text_frame.text`; aquí se
clonan párrafos de la plantilla (con su `pPr`/`rPr`) y se sustituye el `a:t`.
La altura de los textos se estima para paginar y, en los cuadros de tamaño
fijo, reducir la fuente lo justo; el contenido del informe nunca se recorta
ni se reescribe (02_informe.md es lo que se exporta).
"""
from __future__ import annotations

import math
import re
from copy import deepcopy
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from .formato_md import partir_recomendaciones, textos_informe

PLANTILLA = Path(__file__).resolve().parents[1] / "config" / "plantilla_informe.pptx"

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

# Diapositivas de la plantilla (índice 0-based)
(P_PORTADA, P_INDICE, P_INTRO, P_CAB_RESUMEN, P_RESUMEN, P_CAB_DETALLE, P_DETALLE,
 P_DETALLE_ALT, P_CAB_SUGERENCIAS, P_SUGERENCIA, P_CAB_ANEXO) = range(11)

# Banda «RIESGO …» de la tabla de conclusiones: Alto y Medio son los de la plantilla
# (d7 y d8); Crítico y Bajo se derivan de su paleta (línea de acento y gris claro).
COLOR_BANDA = {"Crítico": "958473", "Alto": "6FB1AF", "Medio": "A6A6A6", "Bajo": "BFBFBF"}
GRIS_DETALLES = "F2F2F2"        # bg1 al 95 %: la caja de «detalles descriptivos» de la plantilla
BORDE_BLOQUE_INTRO = "AC9F92"   # marco de los bloques de la introducción
ESTILO_TABLA = "{5940675A-B579-460E-94D1-54222C63F5DA}"   # «Sin estilo, cuadrícula» (el de la plantilla)

# Estimación de altura de texto (ABC Monument Grotesk): anchura media de carácter
# en em e interlineado. Solo se usa para paginar y para ajustar la fuente en
# cuadros de tamaño fijo; el texto del informe no se toca.
FACTOR_ANCHO = 0.47
INTERLINEADO = 1.18
FUENTE_MINIMA = 8.0

SECCIONES = ("Introducción", "Resumen ejecutivo", "Detalle de conclusiones", "Sugerencias de mejora", "Anexo")


# ============================================================ Markdown -> líneas
def _limpiar_md(texto: str) -> str:
    t = re.sub(r"\*\*([^*]+)\*\*", r"\1", texto or "")
    t = re.sub(r"(?<!\w)[*_]([^*_\n]+)[*_](?!\w)", r"\1", t)
    t = re.sub(r"^\s*#+\s*", "", t, flags=re.M)
    t = re.sub(r"^\|?\s*-{3,}\s*(\|\s*-{3,}\s*)*\|?\s*$", "", t, flags=re.M)
    t = re.sub(r"^\|\s*(.*?)\s*\|\s*$", lambda m: " · ".join(x.strip() for x in m.group(1).split("|")), t, flags=re.M)
    return t.replace("‑", "-")


def _lineas_md(texto: str) -> list[str]:
    """Markdown sencillo -> lista de párrafos para la diapositiva. Las líneas de
    un mismo párrafo se unen; cada viñeta (`- `, `* `, `/ `) es un párrafo
    propio con el prefijo «/ » de los informes aprobados."""
    salida: list[str] = []
    buf: list[str] = []

    def cerrar():
        if buf:
            salida.append(" ".join(x.strip() for x in buf))
            buf.clear()

    for linea in _limpiar_md(texto).splitlines():
        if not linea.strip():
            cerrar()
            continue
        m = re.match(r"^\s*(?:[-*•/]|\d+[.)])\s+(.*)$", linea)
        if m:
            cerrar()
            salida.append("/ " + m.group(1).strip())
        else:
            buf.append(linea)
    cerrar()
    return salida


# ============================================================ estimación de altura
def _lineas_estimadas(texto: str, ancho_in: float, pt: float) -> int:
    cpl = max(8, int(ancho_in * 72 / (pt * FACTOR_ANCHO)))
    return sum(max(1, math.ceil(len(l) / cpl)) for l in (texto.split("\n") or [""]))


def _alto_in(texto: str, ancho_in: float, pt: float, extra_pt: float = 0.0) -> float:
    return (_lineas_estimadas(texto, ancho_in, pt) * pt * INTERLINEADO + extra_pt) / 72


# ============================================================ XML: párrafos y runs
def _rpr_de(p):
    r = p.find(A + "r")
    if r is not None and r.find(A + "rPr") is not None:
        return deepcopy(r.find(A + "rPr"))
    fin = p.find(A + "endParaRPr")
    if fin is not None:
        rpr = deepcopy(fin)
        rpr.tag = A + "rPr"
        return rpr
    return None


def _clonar_p(proto, texto: str):
    """Copia un párrafo de la plantilla (pPr, rPr del primer run) con otro texto."""
    p = deepcopy(proto)
    for hijo in list(p):
        if hijo.tag in (A + "r", A + "br", A + "fld"):
            p.remove(hijo)
    r = etree.Element(A + "r")
    rpr = _rpr_de(proto)
    if rpr is not None:
        r.append(rpr)
    t = etree.SubElement(r, A + "t")
    t.text = texto
    fin = p.find(A + "endParaRPr")
    if fin is not None:
        fin.addprevious(r)
    else:
        p.append(r)
    return p


def _poner_parrafos(txBody, parrafos: list) -> None:
    for p in txBody.findall(A + "p"):
        txBody.remove(p)
    for p in parrafos:
        txBody.append(p)


def _set_texto_p(p, texto: str) -> None:
    """Sustituye el texto de un párrafo existente conservando el formato del
    primer run (elimina los demás runs)."""
    runs = p.findall(A + "r")
    if not runs:
        nuevo = _clonar_p(p, texto)
        p.getparent().replace(p, nuevo)
        return
    runs[0].find(A + "t").text = texto
    for r in runs[1:]:
        p.remove(r)
    for br in p.findall(A + "br"):
        p.remove(br)


def _txBody(forma):
    return forma.text_frame._txBody


def _escalar_fuente(el, factor: float, minimo: float = FUENTE_MINIMA) -> None:
    if factor >= 1:
        return
    for tag in ("rPr", "endParaRPr", "defRPr"):
        for e in el.iter(A + tag):
            sz = e.get("sz")
            if sz:
                e.set("sz", str(max(int(minimo * 100), int(int(sz) * factor))))


def _pt_de(p, defecto: float) -> float:
    rpr = _rpr_de(p)
    return int(rpr.get("sz")) / 100 if rpr is not None and rpr.get("sz") else defecto


def _no_fill(ln) -> None:
    for hijo in list(ln):
        if hijo.tag in (A + "noFill", A + "solidFill", A + "gradFill", A + "pattFill"):
            ln.remove(hijo)
    ln.insert(0, etree.Element(A + "noFill"))


def _solid_fill(rgb: str):
    fill = etree.Element(A + "solidFill")
    etree.SubElement(fill, A + "srgbClr").set("val", rgb)
    return fill


def _fondo_y_borde(sp, rgb_fondo: str | None, rgb_borde: str | None, ancho_emu: int = 12700) -> None:
    """Relleno sólido y línea de un `p:sp` (sustituye los existentes)."""
    spPr = sp.find(P + "spPr")
    for tag in ("solidFill", "noFill", "gradFill", "ln"):
        for e in spPr.findall(A + tag):
            spPr.remove(e)
    geom = spPr.find(A + "prstGeom")
    pos = list(spPr).index(geom) + 1 if geom is not None else len(spPr)
    if rgb_fondo:
        fill = etree.Element(A + "solidFill")
        etree.SubElement(fill, A + "schemeClr").set("val", "bg1") if rgb_fondo == "bg1" else \
            etree.SubElement(fill, A + "srgbClr").set("val", rgb_fondo)
        spPr.insert(pos, fill)
        pos += 1
    if rgb_borde:
        ln = etree.Element(A + "ln")
        ln.set("w", str(ancho_emu))
        ln.append(_solid_fill(rgb_borde))
        spPr.insert(pos, ln)


# ============================================================ diapositivas
def _rid_de(prs, slide) -> str:
    return next(rid for rid, rel in prs.part.rels.items() if rel.target_part is slide.part)


def _duplicar(prs, origen):
    """Copia una diapositiva: mismas formas (XML) e imágenes; el diseño se comparte."""
    nueva = prs.slides.add_slide(origen.slide_layout)
    for sh in list(nueva.shapes):
        sh._element.getparent().remove(sh._element)
    mapa = {}
    for rid, rel in origen.part.rels.items():
        if rel.reltype.endswith("/slideLayout"):
            continue
        if rel.is_external:
            mapa[rid] = nueva.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            mapa[rid] = nueva.part.rels.get_or_add(rel.reltype, rel.target_part)
    destino = nueva.shapes._spTree
    for el in origen.shapes._spTree:
        if el.tag in (P + "nvGrpSpPr", P + "grpSpPr"):
            continue
        copia = deepcopy(el)
        for e in copia.iter():
            for k, v in list(e.attrib.items()):
                if k.startswith(R) and v in mapa:
                    e.set(k, mapa[v])
        destino.append(copia)
    return nueva


def _reordenar(prs, orden: list, eliminar: list) -> None:
    lst = prs.slides._sldIdLst
    for slide in eliminar:
        rid = _rid_de(prs, slide)
        for sldId in list(lst):
            if sldId.rId == rid:
                lst.remove(sldId)
        prs.part.drop_rel(rid)
    por_rid = {sldId.rId: sldId for sldId in list(lst)}
    for sldId in list(lst):
        lst.remove(sldId)
    for slide in orden:
        lst.append(por_rid[_rid_de(prs, slide)])


def _forma(slide, nombre: str | None = None, top: float | None = None, contiene: str | None = None,
           tipo: int | None = None):
    for sh in slide.shapes:
        if nombre and sh.name != nombre:
            continue
        if top is not None and abs(sh.top / 914400 - top) > 0.15:
            continue
        if contiene and not (sh.has_text_frame and contiene in sh.text_frame.text):
            continue
        if tipo is not None and sh.shape_type != tipo:
            continue
        return sh
    raise KeyError(f"forma no encontrada en la plantilla: {nombre or contiene or top}")


def _quitar(shape) -> None:
    shape._element.getparent().remove(shape._element)


def _tabla(slide):
    return next(sh for sh in slide.shapes if getattr(sh, "has_table", False) and sh.has_table)


def _es_caja_gris(sh) -> bool:
    sp = sh._element
    return sp.tag == P + "sp" and not (sh.has_text_frame and sh.text_frame.text.strip()) and \
        b'lumMod val="95000"' in etree.tostring(sp.find(P + "spPr"))


# ============================================================ constructor
class ConstructorPlantilla:
    def __init__(self, plantilla: str | Path = PLANTILLA):
        self.prs = Presentation(str(plantilla))
        self.textos = textos_informe()
        self.tpl = list(self.prs.slides)
        if len(self.tpl) != 11:
            raise ValueError(f"La plantilla debe tener 11 diapositivas y tiene {len(self.tpl)}")
        self.protos = self._prototipos()
        self.avisos: list[str] = []

    # ------------------------------------------------ prototipos de párrafo
    def _prototipos(self) -> dict:
        """Párrafos de las tablas de la plantilla que se clonan para el contenido."""
        def celdas(slide):
            return [c.text_frame.paragraphs for c in _tabla(slide).table.rows[0].cells]

        d7, d10 = celdas(self.tpl[P_DETALLE]), celdas(self.tpl[P_SUGERENCIA])
        pe = lambda ps, i: deepcopy(ps[i]._p)  # noqa: E731
        return {
            "rec": {"titulo": pe(d7[1], 0), "blanco": pe(d7[1], 1), "cuerpo": pe(d7[1], 2),
                    "det_intro": pe(d7[1], 5), "det_item": pe(d7[1], 6),
                    "rec_titulo": pe(d7[2], 0), "rec_hueco": pe(d7[2], 1), "rec_cuerpo": pe(d7[2], 2), "rec_blanco": pe(d7[2], 3),
                    "etiqueta": pe(d7[3], 0), "valor": pe(d7[3], 1), "meta_blanco": pe(d7[3], 2)},
            "sug": {"titulo": pe(d10[1], 0), "blanco": pe(d10[1], 1), "cuerpo": pe(d10[1], 2),
                    "det_intro": pe(d10[1], 4), "det_item": pe(d10[1], 5),
                    "rec_titulo": pe(d10[2], 0), "rec_hueco": pe(d10[2], 1), "rec_cuerpo": pe(d10[2], 2), "rec_blanco": pe(d10[2], 3),
                    "etiqueta": pe(d10[3], 0), "valor": pe(d10[3], 2), "meta_blanco": pe(d10[3], 1)},
        }

    # ------------------------------------------------ 1. portada
    def portada(self, proyecto: dict):
        s = self.tpl[P_PORTADA]
        _set_texto_p(_forma(s, "object 2", top=3.09).text_frame.paragraphs[0]._p, proyecto.get("nombre", ""))
        fecha = _forma(s, "object 2", top=6.88)
        ps = fecha.text_frame.paragraphs
        _set_texto_p(ps[0]._p, proyecto.get("fecha", ""))
        if len(ps) > 1:
            _set_texto_p(ps[1]._p, f"Ref.: {proyecto.get('referencia', '')}")
        lista = _forma(s, "object 4")
        proto = lista.text_frame.paragraphs[0]._p
        _poner_parrafos(_txBody(lista), [_clonar_p(proto, x) for x in (proyecto.get("distribucion") or [""])])
        return s

    # ------------------------------------------------ 3. introducción
    def _estructura_intro(self, texto: str, anio: str) -> tuple[str, str, list[tuple[str, list[str]]]]:
        plan = self.textos["plan_auditoria"].format(anio=anio)
        normas = self.textos["normas"]
        bloques: list[tuple[str, list[str]]] = []
        libre: list[str] = []
        actual: list[str] | None = None
        for parrafo in re.split(r"\n\s*\n", (texto or "").strip()):
            lineas = [l for l in parrafo.splitlines() if l.strip()]
            if not lineas:
                continue
            primera = _limpiar_md(lineas[0]).strip()
            if primera.startswith("La auditoría ha sido realizada en cumplimiento"):
                plan = primera
                continue
            if primera.startswith("El trabajo ha sido llevado a cabo"):
                normas = primera
                continue
            # etiqueta de bloque: «**Contexto:** …» (negrita) o una de las etiquetas conocidas
            m = re.match(r"^\*\*([^*\n]{2,60}?):?\*\*:?\s*(.*)$", lineas[0].strip()) or \
                re.match(r"^(Contexto|Objetivo(?: de la auditoría)?|Riesgos a cubrir|Alcance(?: de la auditoría)?|"
                         r"Principales magnitudes|Magnitudes)\s*:\s*(.*)$", primera)
            if m:
                actual = []
                bloques.append((m.group(1).strip().rstrip(":"), actual))
                lineas = ([m.group(2)] if m.group(2).strip() else []) + lineas[1:]
            destino = actual if actual is not None else libre
            destino.extend(_lineas_md("\n".join(lineas)))
        if libre:
            bloques.insert(0, ("", libre))
        return plan, normas, bloques

    def introduccion(self, texto: str, anio: str):
        s = self.tpl[P_INTRO]
        plan, normas, bloques = self._estructura_intro(texto, anio)
        _set_texto_p(_forma(s, contiene="La auditoría ha sido realizada").text_frame.paragraphs[0]._p, plan)
        _set_texto_p(_forma(s, "CuadroTexto 45").text_frame.paragraphs[0]._p, normas)
        cuadros = [sh for sh in s.shapes if sh._element.tag == P + "sp" and sh.has_text_frame
                   and re.match(r"^(Contexto|Objetivo|Riesgos|Alcance)", sh.text_frame.text)]
        marcos = [sh for sh in s.shapes if sh._element.tag == P + "sp" and sh.has_text_frame and sh.text_frame.text.strip() == "7"]
        for m in marcos:                      # el marco pasa a ser el borde del propio cuadro de texto
            _quitar(m)
        proto_p = deepcopy(cuadros[0].text_frame.paragraphs[0]._p)
        proto_sp = deepcopy(cuadros[0]._element)
        while len(cuadros) < len(bloques):
            nuevo = deepcopy(proto_sp)
            nuevo.find(".//" + P + "cNvPr").set("id", str(1000 + len(cuadros)))
            s.shapes._spTree.append(nuevo)
            cuadros.append(list(s.shapes)[-1])
        for sobrante in cuadros[len(bloques):]:
            _quitar(sobrante)
        cuadros = cuadros[:len(bloques)]

        columnas = [(0.75, 6.42), (7.45, 5.45)]
        y0, y_max, hueco = 1.80, 6.85, 0.12
        pt = _pt_de(proto_p, 11.0)
        while True:
            posiciones, ok = self._distribuir_bloques(bloques, columnas, y0, y_max, hueco, pt)
            if ok or pt <= FUENTE_MINIMA:
                break
            pt -= 0.5
        if not ok:
            self.avisos.append("Introducción: el texto excede la diapositiva incluso con fuente mínima.")
        for cuadro, (etiqueta, entradas), (x, y, w, h) in zip(cuadros, bloques, posiciones):
            parrafos = []
            entradas = list(entradas) or [""]
            if etiqueta and not entradas[0].startswith("/ "):
                entradas[0] = f"{etiqueta}: {entradas[0]}"
            elif etiqueta:
                entradas.insert(0, f"{etiqueta}:")
            for e in entradas:
                p = _clonar_p(proto_p, e)
                if e.startswith("/ "):
                    for spc in p.iter(A + "spcAft"):
                        for pts in spc.iter(A + "spcPts"):
                            pts.set("val", "200")
                parrafos.append(p)
            _poner_parrafos(_txBody(cuadro), parrafos)
            _escalar_fuente(_txBody(cuadro), pt / _pt_de(proto_p, 11.0))
            _fondo_y_borde(cuadro._element, "bg1", BORDE_BLOQUE_INTRO)
            cuadro.left, cuadro.top, cuadro.width, cuadro.height = Inches(x), Inches(y), Inches(w), Inches(h)
        return s

    @staticmethod
    def _distribuir_bloques(bloques, columnas, y0, y_max, hueco, pt):
        posiciones, col, y, ok = [], 0, y0, True
        for etiqueta, entradas in bloques:
            texto = "\n".join(([etiqueta + ": "] if etiqueta else []) + list(entradas))
            x, w = columnas[col]
            h = _alto_in(texto, w - 0.2, pt, extra_pt=4 * max(1, len(entradas))) + 0.16
            if y + h > y_max and y > y0:
                if col + 1 < len(columnas):
                    col, y = col + 1, y0
                    x, w = columnas[col]
                    h = _alto_in(texto, w - 0.2, pt, extra_pt=4 * max(1, len(entradas))) + 0.16
                else:
                    ok = False
            if y + h > y_max:
                ok = False
            posiciones.append((x, y, w, h))
            y += h + hueco
        return posiciones, ok

    # ------------------------------------------------ 4. resumen ejecutivo
    def resumen(self, texto: str, evaluacion: str):
        s = self.tpl[P_RESUMEN]
        cuadro = _forma(s, "Rectángulo 9")
        proto = cuadro.text_frame.paragraphs[0]._p
        lineas = _lineas_md(texto) or [""]
        _poner_parrafos(_txBody(cuadro), [_clonar_p(proto, l) for l in lineas])
        pt = _pt_de(proto, 12.0)
        alto = _alto_in("\n".join(lineas), cuadro.width / 914400 - 0.2, pt, extra_pt=5.45 * len(lineas))
        disponible = cuadro.height / 914400 - 0.1
        if alto > disponible:
            _escalar_fuente(_txBody(cuadro), max(FUENTE_MINIMA / pt, disponible / alto))
        # punto de la escala de Evaluación Global (posición calibrada sobre la plantilla)
        escala = [x.lower() for x in self.textos["escala_evaluacion_global"]]
        punto = _forma(s, "Elipse 8")
        if (evaluacion or "").strip().lower() in escala:
            k = escala.index(evaluacion.strip().lower())
            punto.top = Inches(4.558 + (k - 3) * 0.2558)
        else:
            _quitar(punto)
            if evaluacion:
                self.avisos.append(f"Evaluación global «{evaluacion}» no está en la escala; se omite el marcador.")
        proximos = _forma(s, "CuadroTexto 37")
        _poner_parrafos(_txBody(proximos), [_clonar_p(proximos.text_frame.paragraphs[0]._p, self.textos["proximos_pasos"])])
        return s

    # ------------------------------------------------ 5/6. detalle (tabla)
    def detalle(self, num: int, c: dict, es_sugerencia: bool) -> list:
        base = self.tpl[P_SUGERENCIA if es_sugerencia else P_DETALLE]
        pr = self.protos["sug" if es_sugerencia else "rec"]
        gf0 = _tabla(base)
        anchos = [int(g.get("w")) / 914400 for g in gf0._element.find(".//" + A + "tblGrid")]
        w_cuerpo = anchos[1] - 0.16 - 0.2
        w_rec, w_meta = anchos[2] - 0.2, anchos[3] - 0.16
        pt_cuerpo = _pt_de(pr["cuerpo"], 11.0)
        pt_titulo = _pt_de(pr["titulo"], 16.0)

        titulo = f"{num:02d} {(c.get('titulo') or '').strip()}"
        elementos: list[tuple[str, str, float]] = []      # (tipo, texto, alto_in)

        def add(tipo, texto, pt=pt_cuerpo, ancho=w_cuerpo, extra=6.0):
            elementos.append((tipo, texto, _alto_in(texto, ancho, pt, extra_pt=extra)))

        add("titulo", titulo, pt_titulo, extra=2)
        add("blanco", "", extra=0)
        for par in _lineas_md(c.get("incidencia", "")) + _lineas_md(c.get("causa_raiz", "")):
            add("cuerpo", par)
        items = [re.sub(r"^\s*[-*•/]\s*", "", l).strip() for l in (c.get("como_se_ha_llegado") or "").splitlines() if l.strip()]
        if items:
            add("det_intro", self.textos["detalles_descriptivos"], extra=2)
            for it in items:
                add("det_item", "/ " + _limpiar_md(it), ancho=w_cuerpo - 0.2, extra=2)
        for par in _lineas_md(c.get("consecuencias", "")):
            add("cons", par)

        alto_tabla0 = int(gf0._element.find(".//" + A + "tr").get("h")) / 914400
        cap_primera, cap_resto = alto_tabla0 - 0.3, 5.89 - 0.3
        paginas = self._paginar(elementos, cap_primera if (not es_sugerencia or num == 1) else cap_resto, cap_resto)

        etiqueta_rec = "Sugerencia de mejora" if es_sugerencia else "Recomendación"
        recs = partir_recomendaciones(c.get("recomendacion", ""))
        slides = []
        for k, pagina in enumerate(paginas):
            s = _duplicar(self.prs, base)
            for sh in [x for x in s.shapes if _es_caja_gris(x)]:
                _quitar(sh)                                    # la caja gris pasa a ser una fila de la tabla
            gf = _tabla(s)
            if es_sugerencia and (num > 1 or k > 0):           # solo la primera sugerencia lleva el párrafo fijo
                _quitar(_forma(s, contiene="dado su impacto limitado"))
                gf.top, gf.height = Inches(1.01), Inches(5.89)
            alto_tabla = gf.height / 914400
            grupos = self._grupos(pagina, pr, k, titulo)
            self._rellenar_tabla(gf, grupos, alto_tabla, c, es_sugerencia, pr, recs, etiqueta_rec, num, w_rec, w_meta)
            slides.append(s)
        return slides

    @staticmethod
    def _paginar(elementos, cap_primera: float, cap_resto: float) -> list[list]:
        """Reparte los elementos del cuerpo en páginas por altura estimada. La
        línea «A continuación…» va siempre con al menos un detalle."""
        paginas: list[list] = [[]]
        cap, usado = cap_primera, 0.0
        i = 0
        while i < len(elementos):
            tipo, texto, h = elementos[i]
            unidad = [elementos[i]]
            if tipo == "det_intro" and i + 1 < len(elementos):
                unidad.append(elementos[i + 1])
            h_unidad = sum(x[2] for x in unidad)
            if usado + h_unidad > cap and any(t not in ("titulo", "blanco") for t, _, _ in paginas[-1]):
                paginas.append([])
                cap, usado = cap_resto, 0.45   # título de continuación + línea en blanco
            paginas[-1].extend(unidad)
            usado += h_unidad
            i += len(unidad)
        return paginas

    @staticmethod
    def _grupos(pagina, pr, k, titulo) -> list[tuple[str, list, float]]:
        """Elementos de una página -> filas de la tabla: [(tipo_fila, párrafos, alto_estimado)]."""
        filas: list[tuple[str, list, float]] = []
        if k > 0:
            filas.append(["cuerpo", [_clonar_p(pr["titulo"], f"{titulo} (continuación)"), _clonar_p(pr["blanco"], "")], 0.45])
        for tipo, texto, h in pagina:
            fila = {"titulo": "cuerpo", "blanco": "cuerpo", "cuerpo": "cuerpo", "det_intro": "detalles",
                    "det_item": "detalles", "cons": "consecuencias"}[tipo]
            if not filas or filas[-1][0] != fila:
                filas.append([fila, [], 0.0])
            proto = {"titulo": "titulo", "blanco": "blanco", "cuerpo": "cuerpo", "det_intro": "det_intro",
                     "det_item": "det_item", "cons": "cuerpo"}[tipo]
            filas[-1][1].append(_clonar_p(pr[proto], texto))
            filas[-1][2] += h
        return [tuple(f) for f in filas]

    def _rellenar_tabla(self, gf, grupos, alto_tabla, c, es_sugerencia, pr, recs, etiqueta_rec, num, w_rec, w_meta):
        tbl = gf._element.find(".//" + A + "tbl")
        tr0 = tbl.find(A + "tr")
        n = len(grupos)
        trs = [tr0] + [deepcopy(tr0) for _ in range(n - 1)]
        for tr in trs[1:]:
            tbl.append(tr)
        # alturas: proporcionales a lo estimado, mínimo 0,5", suman la altura de la tabla
        estimados = [max(0.5, g[2] + 0.2) for g in grupos]
        factor = alto_tabla / sum(estimados) if sum(estimados) > alto_tabla else 1.0
        alturas = [h * factor for h in estimados]
        alturas[-1] += alto_tabla - sum(alturas)   # el sobrante queda abajo: contenido compacto arriba
        for i, (tr, (tipo, parrafos, _)) in enumerate(zip(trs, grupos)):
            tr.set("h", str(int(Inches(alturas[i]))))
            tcs = tr.findall(A + "tc")
            _poner_parrafos(tcs[1].find(A + "txBody"), parrafos)
            tcPr = tcs[1].find(A + "tcPr")
            if tipo == "detalles":
                for e in tcPr.findall(A + "solidFill") + tcPr.findall(A + "noFill"):
                    tcPr.remove(e)
                tcPr.append(_solid_fill(GRIS_DETALLES))
            for j, tc in enumerate(tcs):
                pr_tc = tc.find(A + "tcPr")
                if i > 0 and pr_tc.find(A + "lnT") is not None:
                    _no_fill(pr_tc.find(A + "lnT"))
                if i < n - 1 and pr_tc.find(A + "lnB") is not None:
                    _no_fill(pr_tc.find(A + "lnB"))
                if j != 1 and n > 1:
                    if i == 0:
                        tc.set("rowSpan", str(n))
                    else:
                        tc.set("vMerge", "1")
                        tb = tc.find(A + "txBody")
                        _poner_parrafos(tb, [_clonar_p(pr["blanco"], "")])
        celdas = tr0.findall(A + "tc")
        # banda RIESGO
        nivel = (c.get("nivel_riesgo") or "").strip().capitalize()
        tb0 = celdas[0].find(A + "txBody")
        p_banda = next(p for p in tb0.findall(A + "p") if p.findall(A + "r"))
        _set_texto_p(p_banda, f"RIESGO {nivel.upper()}".strip())
        _poner_parrafos(tb0, [p_banda])   # sin párrafos vacíos: en texto vertical serían columnas
        if not es_sugerencia:
            tcPr0 = celdas[0].find(A + "tcPr")
            for e in tcPr0.findall(A + "solidFill") + tcPr0.findall(A + "noFill"):
                tcPr0.remove(e)
            tcPr0.append(_solid_fill(COLOR_BANDA.get(nivel, "A6A6A6")))
        # recomendaciones
        ps, texto_rec = [], []
        for k, rec in enumerate(recs, 1):
            if ps:
                ps.append(_clonar_p(pr["rec_blanco"], ""))
            ps += [_clonar_p(pr["rec_titulo"], f"{etiqueta_rec} {num}.{k}"), _clonar_p(pr["rec_hueco"], "")]
            ps += [_clonar_p(pr["rec_cuerpo"], l) for l in _lineas_md(rec)]
            texto_rec += [f"{etiqueta_rec} {num}.{k}"] + _lineas_md(rec)
        if c.get("referencia_recomendacion"):
            ps += [_clonar_p(pr["rec_blanco"], ""), _clonar_p(pr["rec_cuerpo"], f"Ref.-{c['referencia_recomendacion']}")]
            texto_rec.append(f"Ref.-{c['referencia_recomendacion']}")
        if not ps:
            ps = [_clonar_p(pr["rec_cuerpo"], "")]
        tb_rec = celdas[2].find(A + "txBody")
        _poner_parrafos(tb_rec, ps)
        pt_rec = _pt_de(pr["rec_cuerpo"], 11.0)
        alto_rec = _alto_in("\n".join(texto_rec), w_rec, pt_rec, extra_pt=8 * max(1, len(recs)))
        if alto_rec > alto_tabla - 0.3:
            _escalar_fuente(tb_rec, (alto_tabla - 0.3) / alto_rec)
        # área / responsable / plazo
        ps = []
        campos = [("Área", "area")] if es_sugerencia else [("Área", "area"), ("Responsable", "responsable"), ("Plazo", "plazo")]
        for idx, (etiqueta, clave) in enumerate(campos):
            if idx:
                ps.append(_clonar_p(pr["meta_blanco"], ""))
            ps.append(_clonar_p(pr["etiqueta"], etiqueta))
            if es_sugerencia:
                ps.append(_clonar_p(pr["meta_blanco"], ""))
            ps += [_clonar_p(pr["valor"], l) for l in (_lineas_md(c.get(clave) or "") or ["Pendiente"])]
        _poner_parrafos(celdas[3].find(A + "txBody"), ps)

    # ------------------------------------------------ 7. anexo de planes de acción
    def anexo(self, conclusiones: list[dict]) -> list:
        filas = []
        for i, c in enumerate(conclusiones, 1):
            recs = partir_recomendaciones(c.get("recomendacion", "")) or [""]
            for k, _ in enumerate(recs, 1):
                filas.append([f"{i:02d}" if k == 1 else "", c.get("titulo", "") if k == 1 else "", f"{i}.{k}", "",
                              " · ".join(x for x in (c.get("responsable", ""), c.get("area", "")) if x), c.get("plazo", "")])
        por_pagina = 8
        slides = []
        for pag in range(0, max(len(filas), 1), por_pagina):
            s = _duplicar(self.prs, self.tpl[P_DETALLE])
            _quitar(_tabla(s))
            for sh in [x for x in s.shapes if _es_caja_gris(x)]:
                _quitar(sh)
            titulo = _forma(s, "object 3")
            _set_texto_p(titulo.text_frame.paragraphs[0]._p, "Anexo: planes de acción")
            titulo.width = Inches(11.5)
            bloque = filas[pag:pag + por_pagina] or [["", "", "", "", "", ""]]
            forma = s.shapes.add_table(len(bloque) + 1, 6, Inches(0.36), Inches(1.05), Inches(12.6), Inches(0.42) * (len(bloque) + 1))
            tabla = forma.table
            tblPr = forma._element.find(".//" + A + "tblPr")
            estilo = tblPr.find(A + "tableStyleId")
            if estilo is None:
                estilo = etree.SubElement(tblPr, A + "tableStyleId")
            estilo.text = ESTILO_TABLA
            for j, w in enumerate((0.6, 4.5, 0.8, 3.2, 2.5, 1.0)):
                tabla.columns[j].width = Inches(w)
            cabeceras = ("N.º", "Observación", "Rec. N.º", "Plan de Acción", "Persona y Área Responsable", "Plazo")
            for i, fila in enumerate([list(cabeceras)] + bloque):
                for j, v in enumerate(fila):
                    celda = tabla.cell(i, j)
                    p = celda.text_frame.paragraphs[0]
                    r = p.add_run()
                    r.text = v
                    r.font.size = Pt(10)
                    r.font.name = "ABC Monument Grotesk" if i == 0 else "ABC Monument Grotesk Thin"
                    r.font.color.rgb = RGBColor(0x23, 0x1F, 0x20)
                    if i == 0:
                        celda.fill.solid()
                        celda.fill.fore_color.rgb = RGBColor.from_string(GRIS_DETALLES)
            slides.append(s)
        return slides

    # ------------------------------------------------ portadillas
    def portadilla(self, indice: int):
        """Portadilla de sección de la plantilla; se unifican los nombres de las
        secciones con los del índice (la plantilla trae variantes)."""
        s = self.tpl[indice]
        cuadro = _forma(s, contiene="Introducción")
        con_texto = [p for p in cuadro.text_frame.paragraphs if p.text.strip()]
        for p, nombre in zip(con_texto, SECCIONES):
            _set_texto_p(p._p, nombre)
        return s

    # ------------------------------------------------ montaje
    def construir(self, datos: dict) -> None:
        proyecto = datos.get("proyecto", {})
        m_anio = re.search(r"\d{4}", proyecto.get("fecha", "") or "")
        anio = m_anio.group(0) if m_anio else "xxxx"
        conclusiones, sugerencias = datos.get("conclusiones", []), datos.get("sugerencias", [])
        orden = [self.portada(proyecto), self.tpl[P_INDICE], self.introduccion(datos.get("introduccion", ""), anio),
                 self.portadilla(P_CAB_RESUMEN), self.resumen(datos.get("resumen_ejecutivo", ""), datos.get("evaluacion_global", "")),
                 self.portadilla(P_CAB_DETALLE)]
        for i, c in enumerate(conclusiones, 1):
            orden += self.detalle(i, c, es_sugerencia=False)
        if sugerencias:
            orden.append(self.portadilla(P_CAB_SUGERENCIAS))
            for i, c in enumerate(sugerencias, 1):
                orden += self.detalle(i, c, es_sugerencia=True)
        orden.append(self.portadilla(P_CAB_ANEXO))
        orden += self.anexo(conclusiones)
        eliminar = [s for s in self.tpl if s not in orden]
        _reordenar(self.prs, orden, eliminar)

    def guardar(self, ruta: str | Path) -> Path:
        ruta = Path(ruta)
        ruta.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(ruta))
        return ruta


def construir_desde_datos(datos: dict, ruta_salida: str | Path, plantilla: str | Path = PLANTILLA) -> Path:
    """Exporta el informe completo (dict de `parsear_informe` + `proyecto`) sobre la
    plantilla corporativa. Devuelve la ruta del .pptx; los avisos de ajuste
    quedan en `construir_desde_datos.avisos`."""
    c = ConstructorPlantilla(plantilla)
    c.construir(datos)
    construir_desde_datos.avisos = c.avisos  # type: ignore[attr-defined]
    return c.guardar(ruta_salida)


construir_desde_datos.avisos = []  # type: ignore[attr-defined]
