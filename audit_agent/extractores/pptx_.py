"""Extractor de PPTX con python-pptx.

Cada diapositiva es una `Page` (como pedía el spec original) y además una
sección de Markdown por derecho propio: se antepone un `Heading` de nivel 1
con el título de la diapositiva (o un rótulo genérico si no tiene), igual
que en DOCX los `Heading` marcan sección. Así la división en diapositivas se
nota en el propio Markdown, no solo en la lista de páginas.

Simplificaciones deliberadas: el orden de las formas dentro de una
diapositiva es el que da python-pptx (orden del XML), no un análisis de
layout por posición — puede no coincidir con el orden de lectura visual en
diapositivas con formas superpuestas. Las formas agrupadas (`GROUP`) no se
recorren recursivamente.
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image as PILImage

from .base import BaseExtractor
from .models import Block, Heading, Image, Page, Paragraph, Quote, Table
from .ocr import ocr_image


class PPTXExtractor(BaseExtractor):
    supported_formats = frozenset({"pptx"})

    def extract_pages(self, path: Path) -> list[Page]:
        presentation = Presentation(str(path))
        return [
            Page(number=index, blocks=_blocks_for_slide(index, slide))
            for index, slide in enumerate(presentation.slides, start=1)
        ]


def _blocks_for_slide(index: int, slide) -> list[Block]:
    blocks: list[Block] = [Heading(text=_slide_title(index, slide), level=1)]

    # `slide.shapes.title` construye un wrapper nuevo en cada acceso, así que
    # comparar por identidad (`is`) con los de `slide.shapes` nunca coincide;
    # `shape_id` es el identificador estable del XML subyacente.
    title_shape = slide.shapes.title
    title_id = title_shape.shape_id if title_shape is not None else None
    for shape in slide.shapes:
        if shape.shape_id == title_id:
            continue
        blocks.extend(_blocks_for_shape(shape))

    notes = _slide_notes(slide)
    if notes:
        blocks.append(Heading(text="Notas del orador", level=2))
        blocks.append(Quote(text=notes))

    return blocks


def _slide_title(index: int, slide) -> str:
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.has_text_frame:
        text = title_shape.text_frame.text.strip()
        if text:
            return text
    return f"Diapositiva {index}"


def _blocks_for_shape(shape) -> list[Block]:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return [_image_block(shape)]
    if shape.has_table:
        return [_table_block(shape.table)]
    if shape.has_text_frame:
        return _paragraphs_for_text_frame(shape.text_frame)
    return []


def _paragraphs_for_text_frame(text_frame) -> list[Paragraph]:
    blocks = []
    for paragraph in text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            blocks.append(Paragraph(text=text))
    return blocks


def _table_block(table) -> Table:
    rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
    headers = rows[0] if rows else []
    body_rows = rows[1:]
    return Table(headers=headers, rows=body_rows)


def _image_block(shape) -> Image:
    with PILImage.open(BytesIO(shape.image.blob)) as pil_image:
        ocr_text = ocr_image(pil_image.convert("RGB"))
    return Image(caption="Imagen incrustada", ocr_text=ocr_text)


def _slide_notes(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()
