"""
Lectores de exportaciones de Pentana (y de cualquier documento de entrada/).

Interfaz común: `leer(ruta) -> Documento` con el texto normalizado a Markdown
listo para el extractor. Cada lector conserva la estructura que exista en el
origen (títulos, apartados, tablas, párrafos) y NUNCA inventa secciones.

Añadir un formato nuevo = escribir una función `_leer_xxx(ruta) -> str` y
registrarla en `LECTORES` por extensión. Cuando se conozca el formato real
de la exportación de Pentana, se añade aquí (p. ej. un lector específico que
reconozca sus cabeceras "Tarea / Prueba realizada / Resultados / Comentarios
del área / Evidencias") sin tocar el resto de la herramienta.

Lo que se envía al modelo queda registrado en trazas/ (ver acciones.extraer)
con el nombre del lector usado, para poder auditar la fidelidad de la lectura.
"""
from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path


class LecturaError(RuntimeError):
    pass


@dataclass
class Documento:
    nombre: str
    lector: str
    texto: str        # Markdown normalizado
    avisos: list[str]


# ------------------------------------------------------------------ utilidades
def _tabla_md(filas: list[list[str]]) -> str:
    filas = [[(c or "").strip().replace("\n", " ").replace("|", "\\|") for c in f] for f in filas]
    filas = [f for f in filas if any(f)]
    if not filas:
        return ""
    ancho = max(len(f) for f in filas)
    filas = [f + [""] * (ancho - len(f)) for f in filas]
    cab, cuerpo = filas[0], filas[1:]
    out = ["| " + " | ".join(cab) + " |", "|" + "---|" * ancho]
    out += ["| " + " | ".join(f) + " |" for f in cuerpo]
    return "\n".join(out)


def _limpiar(texto: str) -> str:
    texto = texto.replace("\r\n", "\n").replace("\r", "\n")
    texto = re.sub(r"[ \t]+\n", "\n", texto)
    texto = re.sub(r"\n{3,}", "\n\n", texto)
    return texto.strip() + "\n"


# ------------------------------------------------------------------ lectores
def _desexcelizar(texto: str) -> str:
    """Normaliza texto pegado desde Excel (caso real de los papeles de trabajo):
    tabulaciones de celdas vacías al final de línea, celdas multilínea entre
    comillas (con `""` como comilla literal) y ráfagas de líneas en blanco."""
    lineas = [l.rstrip("\t ") for l in texto.replace("\r\n", "\n").split("\n")]
    salida: list[str] = []
    en_celda = False
    for l in lineas:
        if not en_celda and l.startswith('"') and not (len(l) > 1 and l.endswith('"') and l.count('"') % 2 == 0):
            en_celda = True
            l = l[1:]
        elif not en_celda and l.startswith('"') and l.endswith('"') and len(l) > 1:
            l = l[1:-1]
        elif en_celda and l.endswith('"') and (len(l) - len(l.rstrip('"'))) % 2 == 1:
            en_celda = False
            l = l[:-1]
        salida.append(l.replace('""', '"'))
    return "\n".join(salida)


def _leer_texto(ruta: Path) -> tuple[str, list[str]]:
    bruto = ruta.read_text(encoding="utf-8", errors="replace")
    avisos = []
    if "\t" in bruto or re.search(r'^"', bruto, re.M):
        bruto = _desexcelizar(bruto)
        avisos.append("Texto pegado desde Excel: se han normalizado tabulaciones y celdas entre comillas.")
    return _limpiar(bruto), avisos


def _leer_docx(ruta: Path) -> tuple[str, list[str]]:
    try:
        import docx
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except ImportError as exc:
        raise LecturaError("Para leer .docx instala python-docx.") from exc
    d = docx.Document(str(ruta))
    partes: list[str] = []
    avisos: list[str] = []
    # Recorre el cuerpo en orden (párrafos y tablas intercalados)
    for el in d.element.body.iterchildren():
        if el.tag == qn("w:p"):
            p = Paragraph(el, d)
            t = p.text.strip()
            if not t:
                continue
            estilo = (p.style.name or "").lower() if p.style is not None else ""
            m = re.match(r"heading (\d)", estilo) or re.match(r"título (\d)", estilo)
            if m:
                partes.append("#" * min(int(m.group(1)) + 1, 6) + " " + t)
            elif estilo.startswith("title") or estilo.startswith("título") and not m:
                partes.append("# " + t)
            elif "list" in estilo or "lista" in estilo:
                partes.append("- " + t)
            else:
                partes.append(t)
        elif el.tag == qn("w:tbl"):
            tabla = Table(el, d)
            filas = [[c.text for c in fila.cells] for fila in tabla.rows]
            partes.append(_tabla_md(filas))
    if not partes:
        avisos.append("El .docx no contiene texto legible.")
    return _limpiar("\n\n".join(partes)), avisos


def _leer_xlsx(ruta: Path) -> tuple[str, list[str]]:
    try:
        import openpyxl
    except ImportError as exc:
        raise LecturaError("Para leer .xlsx instala openpyxl.") from exc
    wb = openpyxl.load_workbook(str(ruta), data_only=True, read_only=True)
    partes: list[str] = []
    avisos: list[str] = []
    for hoja in wb.worksheets:
        filas = [[("" if v is None else str(v)) for v in fila] for fila in hoja.iter_rows(values_only=True)]
        filas = [f for f in filas if any(c.strip() for c in f)]
        if not filas:
            continue
        partes.append(f"## Hoja: {hoja.title}")
        # Bloques separados por filas vacías ya eliminadas: distinguimos
        # "tabla" (>=2 columnas con contenido en la mayoría de filas) de
        # "texto" (una sola celda larga por fila -> párrafo / apartado).
        bloque: list[list[str]] = []

        def volcar():
            if not bloque:
                return
            n_cols = [sum(1 for c in f if c.strip()) for f in bloque]
            if len(bloque) >= 2 and sum(1 for n in n_cols if n >= 2) >= len(bloque) / 2:
                partes.append(_tabla_md(bloque))
            else:
                for f in bloque:
                    celdas = [c.strip() for c in f if c.strip()]
                    if len(celdas) == 1 and len(celdas[0]) < 60 and not celdas[0].endswith("."):
                        partes.append(f"### {celdas[0]}")
                    else:
                        partes.append(" — ".join(celdas))
            bloque.clear()

        for f in filas:
            bloque.append(f)
        volcar()
    wb.close()
    if not partes:
        avisos.append("El .xlsx no contiene celdas con texto.")
    return _limpiar("\n\n".join(partes)), avisos


def _desenvolver(texto: str) -> str:
    """Une las líneas que el PDF parte por el ancho de página: una línea se
    une a la siguiente solo si no termina en puntuación de cierre y la
    siguiente empieza en minúscula (continuación de frase). Los títulos y
    las viñetas quedan como están."""
    lineas = texto.strip().splitlines()
    salida: list[str] = []
    for linea in lineas:
        l = linea.strip()
        if (salida and salida[-1] and not re.search(r"[.:;!?)]$", salida[-1]) and l
                and l[0].islower() and not salida[-1].startswith(("-", "•", "|"))):
            salida[-1] += " " + l
        else:
            salida.append(l)
    return "\n".join(salida)


def _leer_pdf(ruta: Path) -> tuple[str, list[str]]:
    try:
        import pdfplumber
    except ImportError as exc:
        raise LecturaError("Para leer .pdf instala pdfplumber.") from exc
    partes: list[str] = []
    avisos: list[str] = []
    with pdfplumber.open(str(ruta)) as pdf:
        for i, pagina in enumerate(pdf.pages, 1):
            texto = _desenvolver(pagina.extract_text() or "")
            tablas = pagina.extract_tables() or []
            if texto:
                partes.append(texto)
            for t in tablas:
                partes.append(_tabla_md([[c or "" for c in fila] for fila in t]))
            if not texto and not tablas:
                avisos.append(f"Página {i} sin texto extraíble.")
    contenido = _limpiar("\n\n".join(partes))
    if len(contenido.strip()) < 20:
        raise LecturaError(
            f"{ruta.name}: el PDF no tiene capa de texto extraíble (probablemente es un escaneo). "
            "Exporta desde Pentana a Word/Excel/PDF con texto; no se hace OCR.")
    return contenido, avisos


def _leer_pptx(ruta: Path) -> tuple[str, list[str]]:
    """Informes en PowerPoint (p. ej. informes aprobados usados como corpus de
    calibración): una sección `## Diapositiva N` por diapositiva, con el texto
    de cuadros y tablas en orden de lectura (arriba-abajo, izquierda-derecha)."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise LecturaError("Para leer .pptx instala python-pptx.") from exc
    prs = Presentation(str(ruta))
    partes: list[str] = []
    for i, diapositiva in enumerate(prs.slides, 1):
        bloques: list[str] = []
        for sh in sorted(diapositiva.shapes, key=lambda x: ((x.top or 0), (x.left or 0))):
            if sh.has_text_frame:
                texto = "\n".join(p.text.strip() for p in sh.text_frame.paragraphs if p.text.strip())
                if texto:
                    bloques.append(texto)
            elif getattr(sh, "has_table", False) and sh.has_table:
                filas = [[c.text for c in fila.cells] for fila in sh.table.rows]
                if len(filas) > 1 and len(filas[0]) > 1 and all(len(" ".join(f)) < 400 for f in filas):
                    bloques.append(_tabla_md(filas))
                else:  # tablas "de maquetación" (una fila con celdas largas): texto de cada celda
                    for fila in filas:
                        bloques += [c.strip() for c in fila if c.strip()]
        if bloques:
            partes.append(f"## Diapositiva {i}\n\n" + "\n\n".join(bloques))
    if not partes:
        return "", ["El .pptx no contiene texto."]
    return _limpiar("\n\n".join(partes)), []


LECTORES = {
    ".md": ("texto", _leer_texto),
    ".txt": ("texto", _leer_texto),
    ".docx": ("docx", _leer_docx),
    ".xlsx": ("xlsx", _leer_xlsx),
    ".pdf": ("pdf", _leer_pdf),
    ".pptx": ("pptx", _leer_pptx),
}
EXTENSIONES = tuple(LECTORES)


def leer(ruta: str | Path) -> Documento:
    ruta = Path(ruta)
    ext = ruta.suffix.lower()
    if ext not in LECTORES:
        raise LecturaError(f"{ruta.name}: formato no soportado ({ext}). Admitidos: {', '.join(EXTENSIONES)}.")
    nombre_lector, fn = LECTORES[ext]
    texto, avisos = fn(ruta)
    return Documento(nombre=ruta.name, lector=nombre_lector, texto=texto, avisos=avisos)
