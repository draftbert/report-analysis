"""Lectores de entrada/: .md/.txt, .docx (tablas), .xlsx, .pdf con y sin texto."""
from __future__ import annotations

from pathlib import Path

import pytest

from audit_agent.lectores import LecturaError, leer

RAIZ = Path(__file__).resolve().parent.parent
SINTETICOS = RAIZ / "ejemplos" / "entrada_sintetica"
FRASES_CLAVE = ["45 pedidos", "6 de los 45 pedidos", "delegación temporal de permisos", "tres ofertas comparativas",
                "Comentarios del área", "Anexo E2"]


@pytest.mark.parametrize("ext,lector", [("md", "texto"), ("docx", "docx"), ("xlsx", "xlsx"), ("pdf", "pdf")])
def test_lectura_conserva_contenido(ext, lector):
    ruta = (RAIZ / "ejemplos" / "papel_trabajo_compras.md") if ext == "md" else SINTETICOS / f"papel_trabajo_compras.{ext}"
    doc = leer(ruta)
    assert doc.lector == lector and doc.avisos == []
    for frase in FRASES_CLAVE:
        assert frase in doc.texto, (ext, frase)


def test_docx_conserva_estructura_y_tablas():
    texto = leer(SINTETICOS / "papel_trabajo_compras.docx").texto
    assert "## Papel de trabajo" in texto and "#### Resultados" in texto
    assert "| Ref. | Resultado | Evidencia |" in texto and "|---|---|---|" in texto
    assert "| R2 |" in texto
    assert "- Extracto del módulo de compras (Anexo E1)" in texto  # viñeta


def test_xlsx_hojas_como_secciones_y_tabla():
    texto = leer(SINTETICOS / "papel_trabajo_compras.xlsx").texto
    assert "## Hoja: Tarea 3.2" in texto and "## Hoja: Resultados" in texto
    assert "| Ref. | Resultado | Evidencia |" in texto
    assert "### Prueba realizada" in texto


def test_pdf_sin_texto_error_claro(tmp_path):
    vacio = tmp_path / "escaneo.pdf"
    vacio.write_bytes(b"%PDF-1.4\n1 0 obj\n<< /Type /Catalog /Pages 2 0 R >>\nendobj\n2 0 obj\n<< /Type /Pages /Kids [3 0 R] /Count 1 >>\nendobj\n"
                      b"3 0 obj\n<< /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] >>\nendobj\n"
                      b"xref\n0 4\n0000000000 65535 f \n0000000009 00000 n \n0000000058 00000 n \n0000000115 00000 n \n"
                      b"trailer\n<< /Size 4 /Root 1 0 R >>\nstartxref\n190\n%%EOF\n")
    with pytest.raises(LecturaError, match="capa de texto"):
        leer(vacio)


def test_formato_no_soportado(tmp_path):
    f = tmp_path / "x.csv"
    f.write_text("a,b", encoding="utf-8")
    with pytest.raises(LecturaError, match="no soportado"):
        leer(f)


def test_expediente_lee_entrada_con_lectores(expediente_tmp):
    import shutil
    shutil.copy(SINTETICOS / "papel_trabajo_compras.docx", expediente_tmp.ruta / "papeles_trabajo")
    (expediente_tmp.ruta / "contexto" / "design_thinking.md").write_text("# Motivo\nRevisar compras.\n", encoding="utf-8")
    docs = expediente_tmp.leer_entrada()
    assert [(d.carpeta, d.nombre, d.lector) for d in docs] == [
        ("contexto", "design_thinking.md", "texto"),
        ("papeles_trabajo", "papel_trabajo_compras.docx", "docx"), ("papeles_trabajo", "papel_trabajo_compras.md", "texto")]


def test_extraer_registra_lector_y_texto_en_trazas(contexto):
    from audit_agent.acciones import accion_extraer
    from audit_agent.esquemas import ExtraccionConclusiones
    import json
    contexto.llm.respuestas["extraer"] = ExtraccionConclusiones(conclusiones=[])
    accion_extraer(contexto)
    traza = next(contexto.exp.ruta.glob("trazas/*_extraer-entrada.json"))
    d = json.loads(traza.read_text(encoding="utf-8"))
    assert d["documentos"][0]["carpeta"] == "papeles_trabajo" and d["documentos"][0]["lector"] == "texto"
    assert "6 de los 45 pedidos" in d["documentos"][0]["texto_normalizado"]


def test_txt_pegado_desde_excel_se_normaliza():
    doc = leer(RAIZ / "ejemplos" / "papel_trabajo_tarifarios.txt")
    assert doc.lector == "texto" and any("Excel" in a for a in doc.avisos)
    assert "\t" not in doc.texto and "\n\n\n" not in doc.texto
    assert '"zona remota"' in doc.texto                      # "" -> "
    assert "Debilidades del algoritmo CPF\nA raíz" in doc.texto  # comilla de apertura de celda eliminada
    assert "CON INCIDENCIAS" in doc.texto and "TMSCIIF-10" in doc.texto
    assert doc.texto.count('"') % 2 == 0


def test_pptx_texto_por_diapositiva(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.text = "Resumen ejecutivo"
    t = s.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1)).table
    t.cell(0, 0).text, t.cell(0, 1).text, t.cell(1, 0).text, t.cell(1, 1).text = "Área", "Plazo", "Compras", "30/09/2025"
    ruta = tmp_path / "informe.pptx"
    prs.save(str(ruta))
    doc = leer(ruta)
    assert doc.lector == "pptx" and "## Diapositiva 1" in doc.texto
    assert "Resumen ejecutivo" in doc.texto and "| Área | Plazo |" in doc.texto and "| Compras | 30/09/2025 |" in doc.texto


def test_expediente_antiguo_con_entrada_se_lee_como_papeles(tmp_path):
    from audit_agent.expediente import Expediente
    exp = Expediente.crear(tmp_path / "VIEJO", "V", "VIEJO")
    (exp.ruta / "entrada").mkdir()
    (exp.ruta / "entrada" / "pt.md").write_text("Prueba 1. CONCLUSIONES: sin incidencias.\n", encoding="utf-8")
    assert [d.carpeta for d in exp.leer_entrada()] == ["papeles_trabajo"]
