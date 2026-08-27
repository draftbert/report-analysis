"""Exportación del informe sobre la plantilla corporativa (config/plantilla_informe.pptx):
portada, índice, introducción por bloques, portadillas, resumen con marcador de
Evaluación Global, una tabla por recomendación (banda RIESGO, filas prosa /
detalles en gris / consecuencias, Recomendación N.k, Área-Responsable-Plazo),
sugerencias (RIESGO BAJO + párrafo fijo), anexo de planes de acción; sin restos
de la plantilla ni texto «Lorem»."""
from __future__ import annotations

from pptx import Presentation

from audit_agent.ppt_builder import PLANTILLA, _lineas_md, construir_desde_datos

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

CONCLUSION = {
    "titulo": "Limitaciones en el mantenimiento de los tarifarios de última milla", "prueba": "2.11 b)",
    "nivel_riesgo": "Medio", "area": "Transporte", "responsable": "A. Pérez", "plazo": "Fuera de plazo",
    "referencia_recomendacion": "TMSCIIF-10",
    "incidencia": "Las actualizaciones tarifarias se incorporan manualmente a la Herramienta de Costes.",
    "causa_raiz": "La comunicación de los cambios se articula mediante un proceso principalmente manual.",
    "como_se_ha_llegado": "- El volumen de tarifas a parametrizar puede superar las 5.000 combinaciones.\n- Reclamaciones por importe aproximado de 1,4 M€.",
    "consecuencias": "La manualidad del proceso incrementa el riesgo de tarifas desactualizadas.",
    "recomendacion": "Implantar un sistema para la carga y gestión de los tarifarios.\n\nEstablecer una plantilla única de comunicación de acuerdos.",
}
INTRO = ("La auditoría ha sido realizada en cumplimiento del Plan de Auditoría del año 2025, aprobado por la Comisión de Auditoría y Cumplimiento.\n\n"
         "**Contexto:**\nTransporte e-Commerce orquesta los pedidos online.\n\n"
         "**Objetivo de la auditoría:**\nVerificar los controles sobre los tarifarios:\n- Contraste del tarifario negociado.\n- Proceso de gestión del maestro.\n\n"
         "**Riesgos a cubrir:**\n- Error en la asignación del transportista.\n\n"
         "**Alcance de la auditoría:**\nSistemas: TEC, SCA y CPF.\n\n"
         "**Principales magnitudes:**\n- 35M — órdenes.\n\n"
         "El trabajo ha sido llevado a cabo de acuerdo con las Normas Internacionales para la Práctica Profesional de Auditoría Interna, según certificado emitido por el Instituto de Auditores Internos.")
PROYECTO = {"nombre": "Auditoría X", "referencia": "2025_787", "fecha": "Mayo 2025", "distribucion": ["Control de Stock", "Operaciones tienda"]}


def _textos(slide):
    out = []
    for sh in slide.shapes:
        if getattr(sh, "has_table", False) and sh.has_table:
            out += [c.text_frame.text for row in sh.table.rows for c in row.cells]
        elif sh.has_text_frame:
            out.append(sh.text_frame.text)
    return "\n".join(t for t in out if t.strip())


def _tabla(slide):
    return next(sh.table for sh in slide.shapes if getattr(sh, "has_table", False) and sh.has_table)


def _fill(tc) -> str | None:
    f = tc._tc.find(A + "tcPr").find(A + "solidFill")
    return f[0].get("val") if f is not None else None


def _construir(tmp_path, **cambios):
    datos = {"proyecto": PROYECTO, "introduccion": INTRO, "resumen_ejecutivo": "Resumen.\n\n/ Primera conclusión.\n/ Segunda.",
             "evaluacion_global": "Mejorable", "conclusiones": [CONCLUSION] * 8,
             "sugerencias": [{**CONCLUSION, "nivel_riesgo": "Bajo", "recomendacion": "Propuesta X", "referencia_recomendacion": ""}]}
    datos.update(cambios)
    return Presentation(str(construir_desde_datos(datos, tmp_path / "t.pptx")))


def test_plantilla_versionada_y_saneada():
    assert PLANTILLA.exists()
    prs = Presentation(str(PLANTILLA))
    assert len(prs.slides) == 11 and prs.core_properties.author == "" and prs.core_properties.last_modified_by == ""
    assert not any(r.reltype.endswith(("/comments", "/notesSlide", "/tags", "/oleObject")) for s in prs.slides for r in s.part.rels.values())


def test_estructura_completa(tmp_path):
    prs = _construir(tmp_path)
    todo = [_textos(s) for s in prs.slides]
    # portada, índice, intro, portadilla, resumen, portadilla, 8 detalles, portadilla, 1 sugerencia, portadilla, anexo en 2 páginas
    assert len(prs.slides) == 1 + 1 + 1 + 1 + 1 + 1 + 8 + 1 + 1 + 1 + 2
    assert "Lorem" not in "\n".join(todo) and "xxxx" not in "\n".join(todo)
    # portada
    assert "CONFIDENCIAL" in todo[0] and "Informe de Auditoría Interna" in todo[0] and "Auditoría X" in todo[0]
    assert "Control de Stock\nOperaciones tienda" in todo[0] and "Mayo 2025\nRef.: 2025_787" in todo[0]
    # índice y portadillas con los nombres unificados de sección
    assert todo[1].split("\n")[0] == "Introducción" and "Anexo" in todo[1]
    for k in (3, 5, 14, 16):
        assert "Sugerencias de mejora" in todo[k] and "Anexos" not in todo[k] and "Detalle de sugerencias" not in todo[k]
    # introducción: frases fijas y bloques con etiqueta, viñetas «/ »
    intro = todo[2]
    assert "Plan de Auditoría del año 2025" in intro and "Normas Internacionales" in intro
    assert "Contexto: Transporte e-Commerce orquesta" in intro and "Objetivo de la auditoría: Verificar" in intro
    assert "Riesgos a cubrir:\n/ Error en la asignación" in intro and "Principales magnitudes:\n/ 35M" in intro
    assert "**" not in intro and "- Contraste" not in intro
    # resumen: texto, marcador de la escala en «Mejorable» (tercer nivel) y próximos pasos fijos
    resumen = prs.slides[4]
    assert "/ Primera conclusión." in _textos(resumen) and "Los destinatarios del presente informe" in _textos(resumen)
    punto = next(sh for sh in resumen.shapes if sh.name == "Elipse 8")
    assert abs(punto.top / 914400 - (4.558 - 0.2558)) < 0.01


def test_diapositiva_de_recomendacion(tmp_path):
    prs = _construir(tmp_path)
    octava = prs.slides[6 + 7]
    assert _textos(octava).startswith("Detalle de conclusiones") or "Detalle de conclusiones" in _textos(octava)
    t = _tabla(octava)
    assert len(t.rows) == 3 and len(t.columns) == 4
    assert t.cell(0, 0).text == "RIESGO MEDIO" and _fill(t.cell(0, 0)) == "A6A6A6"
    assert t.cell(0, 0)._tc.get("rowSpan") == "3" and t.cell(1, 0)._tc.get("vMerge") == "1"
    cuerpo = t.cell(0, 1).text
    assert cuerpo.startswith("08 Limitaciones en el mantenimiento") and "manualmente a la Herramienta" in cuerpo and "proceso principalmente manual" in cuerpo
    detalles = t.cell(1, 1).text
    assert detalles.startswith("A continuación, se muestran los detalles descriptivos") and "/ El volumen de tarifas" in detalles and "/ Reclamaciones" in detalles
    assert _fill(t.cell(1, 1)) == "F2F2F2" and "- El volumen" not in detalles
    assert t.cell(2, 1).text.startswith("La manualidad del proceso")
    recs = t.cell(0, 2).text
    assert "Recomendación 8.1" in recs and "Implantar un sistema" in recs and "Recomendación 8.2" in recs and "Ref.-TMSCIIF-10" in recs
    assert t.cell(0, 3).text.split("\n") == ["Área", "Transporte", "", "Responsable", "A. Pérez", "", "Plazo", "Fuera de plazo"]
    # la caja gris flotante de la plantilla no sobrevive (los detalles son una fila)
    assert not any(sh.name == "Rectángulo 4" for sh in octava.shapes)
    # formato heredado de la plantilla: título 16 pt, cuerpo «ABC Monument Grotesk Thin» 11 pt
    p0, p2 = t.cell(0, 1).text_frame.paragraphs[0], t.cell(0, 1).text_frame.paragraphs[2]
    assert p0.runs[0].font.size.pt == 16 and p2.runs[0].font.size.pt == 11 and p2.runs[0].font.name == "ABC Monument Grotesk Thin"


def test_sugerencia_y_anexo(tmp_path):
    prs = _construir(tmp_path)
    sug = prs.slides[15]
    t = _tabla(sug)
    assert t.cell(0, 0).text == "RIESGO BAJO" and "dado su impacto limitado" in _textos(sug)
    assert "Sugerencia de mejora 1.1" in t.cell(0, 2).text and "Propuesta X" in t.cell(0, 2).text
    assert t.cell(0, 3).text.startswith("Área") and "Transporte" in t.cell(0, 3).text
    assert "Responsable" not in t.cell(0, 3).text and "Plazo" not in t.cell(0, 3).text
    anexo = prs.slides[17]
    assert "Anexo: planes de acción" in _textos(anexo)
    ta = _tabla(anexo)
    assert [c.text for c in ta.rows[0].cells] == ["N.º", "Observación", "Rec. N.º", "Plan de Acción", "Persona y Área Responsable", "Plazo"]
    assert [c.text for c in ta.rows[1].cells][:3] == ["01", CONCLUSION["titulo"], "1.1"] and ta.cell(1, 4).text == "A. Pérez · Transporte"
    assert ta.cell(2, 2).text == "1.2" and ta.cell(2, 0).text == "" and len(ta.rows) == 9
    assert len(_tabla(prs.slides[18]).rows) == 9  # 16 recomendaciones -> 8 + 8


def test_continuacion_cuando_no_cabe(tmp_path):
    larga = {**CONCLUSION, "incidencia": "\n\n".join(f"Párrafo {i}. " + "texto " * 60 for i in range(12))}
    prs = _construir(tmp_path, conclusiones=[larga], sugerencias=[])
    detalles = [s for s in prs.slides if "Detalle de conclusiones" in _textos(s) and any(getattr(sh, "has_table", False) for sh in s.shapes)]
    assert len(detalles) >= 2
    assert _tabla(detalles[0]).cell(0, 1).text.startswith("01 Limitaciones")
    assert "(continuación)" in _tabla(detalles[1]).cell(0, 1).text
    # sin sugerencias no hay portadilla ni diapositiva de sugerencias
    assert not any("dado su impacto limitado" in _textos(s) for s in prs.slides)
    assert "Lorem" not in "\n".join(_textos(s) for s in prs.slides)


def test_niveles_y_evaluacion_desconocida(tmp_path):
    prs = _construir(tmp_path, conclusiones=[{**CONCLUSION, "nivel_riesgo": "Alto", "como_se_ha_llegado": "", "consecuencias": ""}],
                     evaluacion_global="Regular", sugerencias=[])
    t = _tabla(prs.slides[6])
    assert t.cell(0, 0).text == "RIESGO ALTO" and _fill(t.cell(0, 0)) == "6FB1AF" and len(t.rows) == 1
    assert not any(sh.name == "Elipse 8" for sh in prs.slides[4].shapes)
    assert any("Regular" in a for a in construir_desde_datos.avisos)


def test_lineas_md():
    assert _lineas_md("Uno\ndos\n\n- tres\n* cuatro\n/ cinco\n\n**seis** _siete_") == ["Uno dos", "/ tres", "/ cuatro", "/ cinco", "seis siete"]
